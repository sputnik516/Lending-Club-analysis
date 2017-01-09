from pptx import Presentation
from pptx.util import Inches


def currency_format(value):
    """Format currency values. Ex: $543,921.94"""
    return "${:,.2f}".format(value)

def save_pptx(prs, pptx_name):
    """Save powerpoint file"""
    prs.save(pptx_name)


def new_pptx(pptx_name, grouped_data):
    """Template file"""
    prs = Presentation('ppt_template.pptx')

    """Summary Slide"""
    summary_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(summary_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Summary'

    tf = body_shape.text_frame
    tf.text = "Total Outstanding Principal: {}".format(currency_format(grouped_data['out_prncp'].sum()))
    p = tf.add_paragraph()
    p.text = 'Total Recovered: {}'.format(currency_format(grouped_data['recoveries'].sum()))
    p = tf.add_paragraph()
    p.text = 'Total Profit: {}'.format(currency_format(grouped_data['profit_loss'].sum()))
    p = tf.add_paragraph()
    p.text = 'Net Return: {}%'.format(round(grouped_data['profit_loss'].sum() /
                                      grouped_data['funded_amnt_inv'].sum(), 4) * 100)

    """Loans by Grade"""
    title_only_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Loans by Grade'

    rows = len(grouped_data) + 1
    cols = 6
    top = Inches(2.5)
    left = Inches(0)
    width = Inches(9.75)
    height = Inches(0.8)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(.75)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.92)
    table.columns[3].width = Inches(1.93)
    table.columns[4].width = Inches(1.8)
    table.columns[5].width = Inches(1.8)

    # write column headings
    table.cell(0, 0).text = 'Grade'
    table.cell(0, 1).text = 'Loan Total'
    table.cell(0, 2).text = 'Payments Rec\'d'
    table.cell(0, 3).text = 'Outstanding Pricipal'
    table.cell(0, 4).text = 'Recovered from Defaults'
    table.cell(0, 5).text = 'Profit / Loss'

    row = 0
    while row < len(grouped_data):
        table.cell(row + 1, 0).text = str(grouped_data.ix[row].name)
        table.cell(row + 1, 1).text = str(currency_format((grouped_data.iloc[row][0])))
        table.cell(row + 1, 2).text = str(currency_format((grouped_data.iloc[row][1])))
        table.cell(row + 1, 3).text = str(currency_format((grouped_data.iloc[row][2])))
        table.cell(row + 1, 4).text = str(currency_format((grouped_data.iloc[row][3])))
        table.cell(row + 1, 5).text = str(currency_format((grouped_data.iloc[row][4])))

        row += 1

    save_pptx(pptx_name)